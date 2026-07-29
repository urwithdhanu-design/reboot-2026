"""
Generate Reboot 2026 platform overview PowerPoint (Canton ledger).
Run: python docs/generate_platform_deck.py
Output: docs/Reboot-2026-Platform-Overview.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "Reboot-2026-Platform-Overview.pptx"

GREEN = RGBColor(0, 106, 77)
GREEN_DARK = RGBColor(0, 77, 56)
GREEN_LIGHT = RGBColor(232, 245, 239)
WHITE = RGBColor(255, 255, 255)
INK = RGBColor(26, 31, 28)
MUTED = RGBColor(90, 99, 95)
AMBER = RGBColor(217, 119, 6)


def set_slide_bg(slide, rgb: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_footer(slide, text: str = "Reboot 2026 Insurance · Canton Platform"):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED


def add_title_slide(prs, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, GREEN_DARK)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(94, 224, 176)
    bar.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(11.5), Inches(1.5))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    sbox = slide.shapes.add_textbox(Inches(0.75), Inches(3.85), Inches(11), Inches(1.2))
    sp = sbox.text_frame.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(20)
    sp.font.color.rgb = RGBColor(200, 230, 215)

    add_footer(slide, "Canton Network · Daml ledger · insure360-83a36")


def add_section_slide(prs, title: str, subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, GREEN)
    tbox = slide.shapes.add_textbox(Inches(0.75), Inches(2.8), Inches(11.5), Inches(1.2))
    p = tbox.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(0.75), Inches(4.0), Inches(11), Inches(0.8))
        sp = sbox.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(18)
        sp.font.color.rgb = RGBColor(220, 245, 235)


def add_content_slide(prs, title: str, bullets: list[str], note: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GREEN
    accent.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(12.2), Inches(0.9))
    tp = tbox.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = GREEN_DARK

    body = slide.shapes.add_textbox(Inches(0.65), Inches(1.35), Inches(12.0), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if line.startswith("## "):
            p.text = line[3:]
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = GREEN
            p.space_before = Pt(10)
        elif line.startswith("→ "):
            p.text = line
            p.font.size = Pt(14)
            p.font.color.rgb = INK
            p.level = 1
        else:
            p.text = line
            p.font.size = Pt(15)
            p.font.color.rgb = INK
            p.space_after = Pt(6)

    if note:
        nbox = slide.shapes.add_textbox(Inches(0.65), Inches(6.35), Inches(12), Inches(0.55))
        np = nbox.text_frame.paragraphs[0]
        np.text = f"Insight: {note}"
        np.font.size = Pt(11)
        np.font.italic = True
        np.font.color.rgb = AMBER

    add_footer(slide)


def add_table_slide(prs, title: str, headers: list[str], rows: list[list[str]]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    tbox = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(12), Inches(0.7))
    tbox.text_frame.paragraphs[0].text = title
    tbox.text_frame.paragraphs[0].font.size = Pt(26)
    tbox.text_frame.paragraphs[0].font.bold = True
    tbox.text_frame.paragraphs[0].font.color.rgb = GREEN_DARK

    cols, row_count = len(headers), len(rows) + 1
    table = slide.shapes.add_table(row_count, cols, Inches(0.55), Inches(1.25), Inches(12.2), Inches(0.45 * row_count)).table

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GREEN_LIGHT
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = INK

    add_footer(slide)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Reboot 2026 Insurance Platform",
        "End-to-end digital insurance · KYC · Wallet · Canton minting · Claims · Parametric automation",
    )

    add_content_slide(
        prs,
        "Agenda",
        [
            "Platform vision & architecture",
            "Customer journey: KYC → Wallet → Policy → Mint",
            "Claims: manual review & parametric automation",
            "GBP fund flows: premium, vendor reserve, claims pool",
            "AI, blockchain & Canton integration",
            "Vendor portal & admin operations",
            "Cloud deployment & demo access",
        ],
    )

    add_content_slide(
        prs,
        "Platform overview",
        [
            "Reboot 2026 — digital insurance distribution on GCP with Canton as the ledger",
            "Customer web (React) + Admin console + 12+ Java/Python microservices",
            "GBP wallets (demo) + policy certificates minted on Canton via Daml templates",
            "Parametric claims (flight delay, telematics) with ≤£500 auto-settle",
            "Firebase Hosting (same-origin /api/*) → Cloud Run microservices",
            "Production GCP project: insure360-83a36",
        ],
        "Wallet-service holds GBP; Canton ledger mirrors policies and settlements for audit.",
    )

    add_table_slide(
        prs,
        "Microservices map (local ports)",
        ["Service", "Port", "Responsibility"],
        [
            ["KYC service", "8081", "Auth, KYC, AI agent, assistant hints"],
            ["Policy service", "8082", "Quotes, payments, policies, vendors"],
            ["Wallet service", "8089", "Customer wallets, vendor reserve, claims pool"],
            ["Claims service", "8085", "Claim workflow & settlement"],
            ["Parametric service", "8086", "Oracle rules, auto-triggers"],
            ["Blockchain orchestrator", "8088", "Canton mint, ledger, claim settle"],
            ["Chatbot assistance", "8090", "Stallion RAG chatbot"],
            ["Canton ledger", "internal", "Daml sandbox / JSON API for policy NFTs"],
        ],
    )

    add_content_slide(
        prs,
        "Architecture layers",
        [
            "## Experience",
            "Customer app (localhost:5174) · Admin console (localhost:5175) · Vendor portal",
            "## Core services",
            "KYC · Wallet · Policy · Claims · Parametric · Payment · Notification",
            "## Ledger & intelligence",
            "Blockchain orchestrator · Canton · Ledger sidecar · Fraud scorer",
            "## AI",
            "KYC auto-agent · Stallion chatbot · Flight oracle · Screen assistant",
            "## Events (optional)",
            "Pub/Sub: customer-events, policy-events, wallet-events",
        ],
    )

    add_section_slide(prs, "Customer journey", "Identity → Wallet → Cover → On-chain certificate")

    add_content_slide(
        prs,
        "1 · KYC flow",
        [
            "Register → Upload ID + selfie → Submit KYC",
            "AI KYC agent (default ON) or manual admin review",
            "Admin approve → pending_consent (not straight to verified)",
            "Customer accepts digitisation consent → verified",
            "Publishes CustomerVerified event (does not auto-create wallet)",
            "→ Gate for wallet creation and Canton policy mint",
        ],
        "Two consent layers: KYC digitisation vs wallet email consent.",
    )

    add_content_slide(
        prs,
        "2 · Wallet linking",
        [
            "Requires kyc_status = verified",
            "Create secure wallet (demo SHA-256 address) or link external 0x address",
            "Status → connected · publishes WalletLinked event",
            "Policy service retries any pending Canton mints",
            "Premium payment via POST /api/payments/wallet debits customer wallet",
            "Claim payout can auto-provision wallet if address known",
        ],
    )

    add_content_slide(
        prs,
        "3 · Policy issuance & minting",
        [
            "Quote → Pay premium (wallet) → Policy ISSUED (off-chain record)",
            "Premium credits vendor reserve (vendor-vitality, vendor-homeshield)",
            "Mint gates: paid + KYC verified + wallet connected + compliance",
            "Orchestrator mints policy NFT on Canton → status MINTED",
            "Post-mint: parametric rules auto-provisioned (travel / motor)",
            "Admin can force mint from Tokenization queue",
        ],
        "Deferred mint: policy can be ISSUED while waiting for wallet.",
    )

    add_section_slide(prs, "Claims & settlement", "Manual review · Parametric automation · GBP payout")

    add_content_slide(
        prs,
        "4 · Manual claim settlement",
        [
            "Customer submits claim → pending_approval",
            "Admin review → approve (re-verify Canton + coverage cap)",
            "Wallet-service: claims-pool debit → customer wallet credit",
            "Coverage consumed on policy · chain settlement (best-effort)",
            "Status → settled",
            "Open queries block approval until resolved",
        ],
    )

    add_content_slide(
        prs,
        "5 · Parametric settlement",
        [
            "Rules created at mint from quote answers (flight, telematics)",
            "Triggers: admin simulate · live flight oracle poll",
            "Threshold matched → parametric claim created",
            "≤ £500: auto approveAndSettle immediately",
            "> £500: joins manual approval queue",
            "Chain record (oracle → pool) is audit-only; GBP moves at credit-claim",
        ],
    )

    add_content_slide(
        prs,
        "6 · Fund transfer (claim success)",
        [
            "## Premium inflow",
            "Customer wallet −£ → Vendor reserve +£ (vendor_premium)",
            "## Pool funding",
            "Vendor reserve → Claims pool (vendor_contribution)",
            "Admin top-up · Demo seed: £100k claims pool, £50k per vendor reserve",
            "## Claim payout",
            "Claims pool −£ → Customer wallet +£ (idempotent per claimId)",
            "Canton ledger mirrors: claims-pool → customer wallet (on-chain settlement record)",
        ],
        "Claims paid from claims-pool — vendors must contribute reserve first.",
    )

    add_table_slide(
        prs,
        "AI & automation touchpoints",
        ["Capability", "Service", "Role"],
        [
            ["KYC AI agent", "kyc-service", "Auto-approve submissions (toggle in admin)"],
            ["Stallion chatbot", "chatbot :8090", "RAG insurance Q&A for customers"],
            ["Screen assistant", "kyc-service", "Contextual hints per screen"],
            ["Flight oracle", "parametric", "Poll delay data; trigger rules"],
            ["Fraud scorer", "orchestrator", "Heuristic on mint/settle txs"],
            ["Deferred mint", "policy-service", "Retry on WalletLinked event"],
            ["Auto-settle", "claims-service", "Parametric claims ≤ £500"],
        ],
    )

    add_content_slide(
        prs,
        "Blockchain & Canton",
        [
            "Policy NFT mint via blockchain-orchestrator → Canton JSON API",
            "Insurance chain: PoA validator, typed ledgers (POLICY, CLAIMS, AUDIT)",
            "Claim settlement recorded on-chain (parallel to wallet GBP)",
            "Admin: Tokenization queue, Chain Monitor, Blockchain Ledger UI",
            "Canton JSON API + ledger sidecar for Daml template wire-up",
            "Document hashes & fraud scoring at transaction ingest",
        ],
    )

    add_content_slide(
        prs,
        "Vendor ecosystem",
        [
            "Partners: Vitality (health-plan) · HomeShield (home-insurance)",
            "Vendor reserve receives customer premiums (vendor_premium tx)",
            "Vendor portal: fund shared claims pool from reserve",
            "Vitality quote UI: /vendors/vitality → /quote/health-plan",
            "Demo logins: vendor.vitality@example.com / VendorDemo123!",
            "Claim settlement debits claims-pool (not vendor reserve directly)",
        ],
    )

    add_content_slide(
        prs,
        "Admin operations console",
        [
            "KYC review queue · Customer registry",
            "Tokenization / mint queue · Compliance controls",
            "Claims approval · Parametric simulation",
            "Wallet ops: claims pool top-up, vendor reserves",
            "Platform flows reference (login page animation)",
            "Chain observability · Reports · Vendor management",
        ],
    )

    add_section_slide(prs, "Deployment", "Cloud Run + Firebase Hosting · GCP insure360-83a36")

    add_table_slide(
        prs,
        "Live & local URLs",
        ["App", "Local", "Production"],
        [
            ["Customer web", "http://localhost:5174", "https://insure360-83a36.web.app"],
            ["Admin console", "http://localhost:5175", "https://insure360-83a36-admin.web.app"],
            ["Vendor portal", "localhost:5175/vendor/portal", "Same admin hosting"],
            ["Admin login", "admin@reboot2026.local", "Reboot2026!Admin (seeded)"],
            ["Deploy APIs", "local-dev.cmd start", "deploy\\deploy-cloud-run.cmd"],
            ["Deploy UIs", "npm run dev", "deploy\\deploy-firebase.cmd"],
        ],
    )

    add_content_slide(
        prs,
        "Cloud deployment checklist",
        [
            "setup-gcp-project.ps1 — APIs, Artifact Registry, Firebase sites",
            "setup-cloud-sql.ps1 — PostgreSQL instance (optional)",
            "setup-pubsub.cmd — Event bus topics (optional)",
            "deploy-cloud-run.cmd — All microservices + Canton ledger to Cloud Run",
            "deploy-firebase.cmd — Customer + admin SPAs",
            "Deploy flags: enable Cloud SQL and Pub/Sub via setup scripts + env",
        ],
    )

    add_content_slide(
        prs,
        "Demo walkthrough (local)",
        [
            "1. Register → KYC → accept consent → verified",
            "2. Create wallet → WalletLinked → pending policies mint",
            "3. Get quote → Pay from wallet → vendor reserve credited",
            "4. Vendor portal → contribute to claims pool",
            "5. File claim OR simulate parametric flight delay (admin)",
            "6. Approve claim → customer wallet credited from claims-pool",
        ],
        "Start: local-dev.cmd setup && local-dev.cmd start",
    )

    add_content_slide(
        prs,
        "Key design insights",
        [
            "Wallet creation is always user-initiated after KYC — never automatic",
            "Policies can be ISSUED off-chain while waiting for wallet + mint",
            "GBP truth in wallet-service; blockchain is audit & compliance mirror",
            "Parametric ≤£500 auto-settle reduces ops load for small claims",
            "Vendors must fund claims-pool before sustainable claim payouts",
            "Same-origin API via Firebase rewrites — no VITE_API_BASE in production",
        ],
    )

    add_title_slide(
        prs,
        "Thank you",
        "docs/PLATFORM-FLOWS.md · docs/HOSTING-ACCESS.md · deploy/README.md",
    )

    prs.save(OUT)
    print(f"Created: {OUT}")


if __name__ == "__main__":
    build()
